"""
翻訳済みテキストを PPTX に書き戻すモジュール

TextSegment のリストに設定された translated_text を
元の PPTX の対応する位置に書き戻す。
スタイル（フォント、サイズ、色など）は可能な限り維持する。
"""

from typing import List
from pptx import Presentation
from pptx.shapes.base import BaseShape

from translator.models import TextSegment


def apply_translations(prs: Presentation, segments: List[TextSegment]) -> Presentation:
    """
    翻訳済みテキストを PPTX に書き戻す

    Args:
        prs: python-pptx の Presentation オブジェクト
        segments: 翻訳済みテキストを含む TextSegment のリスト

    Returns:
        翻訳が適用された Presentation オブジェクト（同じインスタンス）
    """
    for segment in segments:
        # 翻訳テキストがない場合はスキップ
        if segment.translated_text is None:
            continue

        # スライドを取得
        slide = prs.slides[segment.slide_index]

        # 要素タイプに応じて処理を分岐
        if segment.element_type == 'note':
            _apply_to_note(slide, segment)
        elif segment.element_type == 'table':
            _apply_to_table(slide, segment)
        else:
            # title, body, chart など、テキストフレームを持つ要素
            _apply_to_text_frame(slide, segment)

    return prs


def _apply_to_text_frame(slide, segment: TextSegment) -> None:
    """
    テキストフレームに翻訳テキストを書き戻す

    スタイル維持戦略:
    - パラグラフの最初の Run のスタイルを保持
    - 翻訳テキスト全体をその Run に設定
    - 残りの Run は空文字にする

    Args:
        slide: スライドオブジェクト
        segment: 翻訳済みテキストセグメント
    """
    shape = slide.shapes[segment.shape_index]

    # テキストフレームを持たない場合はスキップ
    if not hasattr(shape, 'text_frame') or not shape.has_text_frame:
        return

    text_frame = shape.text_frame

    # パラグラフインデックスが指定されている場合はそのパラグラフのみ更新
    if segment.paragraph_index is not None:
        if segment.paragraph_index < len(text_frame.paragraphs):
            paragraph = text_frame.paragraphs[segment.paragraph_index]
            _update_paragraph_text(paragraph, segment.translated_text)
    else:
        # パラグラフインデックスがない場合は最初のパラグラフを更新
        if len(text_frame.paragraphs) > 0:
            paragraph = text_frame.paragraphs[0]
            _update_paragraph_text(paragraph, segment.translated_text)


def _apply_to_table(slide, segment: TextSegment) -> None:
    """
    テーブルセルに翻訳テキストを書き戻す

    Args:
        slide: スライドオブジェクト
        segment: 翻訳済みテキストセグメント
    """
    shape = slide.shapes[segment.shape_index]

    # テーブルでない場合はスキップ
    if not hasattr(shape, 'has_table') or not shape.has_table:
        return

    table = shape.table

    # セル座標が範囲内かチェック
    if (segment.cell_row is not None and
        segment.cell_col is not None and
        segment.cell_row < len(table.rows) and
        segment.cell_col < len(table.rows[0].cells)):

        cell = table.cell(segment.cell_row, segment.cell_col)

        # セルのテキストフレームの最初のパラグラフを更新
        if cell.text_frame and len(cell.text_frame.paragraphs) > 0:
            paragraph = cell.text_frame.paragraphs[0]
            _update_paragraph_text(paragraph, segment.translated_text)


def _apply_to_note(slide, segment: TextSegment) -> None:
    """
    スピーカーノートに翻訳テキストを書き戻す

    Args:
        slide: スライドオブジェクト
        segment: 翻訳済みテキストセグメント
    """
    if not slide.has_notes_slide:
        return

    notes_slide = slide.notes_slide
    text_frame = notes_slide.notes_text_frame

    # ノートのテキストフレームを全体更新
    # ノートの場合は Run の概念がないため、直接 text 属性を設定
    text_frame.text = segment.translated_text


def _update_paragraph_text(paragraph, new_text: str) -> None:
    """
    パラグラフのテキストを更新し、スタイルを維持する

    戦略:
    - 最初の Run に新しいテキストを設定（スタイルは維持される）
    - 残りの Run は空文字にする（削除はしない。削除するとレイアウトが崩れる可能性がある）

    Args:
        paragraph: python-pptx の Paragraph オブジェクト
        new_text: 新しいテキスト
    """
    runs = paragraph.runs

    if len(runs) == 0:
        # Run がない場合は新規作成
        paragraph.text = new_text
        return

    # 最初の Run に新しいテキストを設定
    runs[0].text = new_text

    # 残りの Run は空文字にする
    for i in range(1, len(runs)):
        runs[i].text = ""
