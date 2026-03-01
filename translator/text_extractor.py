"""
PPTX からテキストセグメントを抽出するモジュール

python-pptx を使用して、スライド内のテキスト要素を走査し、
TextSegment のリストとして返す。
"""

from typing import List
from pptx import Presentation
from pptx.shapes.base import BaseShape
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Inches

from translator.models import TextSegment


def extract_texts(prs: Presentation, translate_notes: bool = True) -> List[TextSegment]:
    """
    PPTX からすべてのテキストセグメントを抽出する

    Args:
        prs: python-pptx の Presentation オブジェクト
        translate_notes: スピーカーノートを抽出するか

    Returns:
        TextSegment のリスト
    """
    segments: List[TextSegment] = []

    for slide_idx, slide in enumerate(prs.slides):
        # 通常のシェイプからテキストを抽出
        for shape_idx, shape in enumerate(slide.shapes):
            segments.extend(_extract_from_shape(shape, slide_idx, shape_idx))

        # スピーカーノートを抽出
        if translate_notes and slide.has_notes_slide:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()
            if notes_text:
                segments.append(TextSegment(
                    original_text=notes_text,
                    slide_index=slide_idx,
                    shape_index=-1,  # ノートは特別な shape_index として -1 を使用
                    element_type='note'
                ))

    return segments


def _extract_from_shape(shape: BaseShape, slide_idx: int, shape_idx: int) -> List[TextSegment]:
    """
    単一のシェイプからテキストセグメントを抽出する

    Args:
        shape: python-pptx の BaseShape オブジェクト
        slide_idx: スライドインデックス
        shape_idx: シェイプインデックス

    Returns:
        TextSegment のリスト
    """
    segments: List[TextSegment] = []

    # テーブル（GraphicFrame オブジェクトのチェックを先に行う）
    if hasattr(shape, 'has_table') and shape.has_table:
        segments.extend(_extract_from_table(shape, slide_idx, shape_idx))

    # テキストフレームを持つシェイプ（タイトル、本文など）
    elif hasattr(shape, 'text_frame') and shape.has_text_frame:
        segments.extend(_extract_from_text_frame(shape, slide_idx, shape_idx))

    # グループシェイプ（再帰的に処理）
    elif hasattr(shape, 'shapes') and shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for sub_idx, sub_shape in enumerate(shape.shapes):
            segments.extend(_extract_from_shape(sub_shape, slide_idx, shape_idx))

    # グラフ（タイトルのみ）
    elif hasattr(shape, 'has_chart') and shape.has_chart:
        if hasattr(shape.chart, 'chart_title') and shape.chart.has_title:
            title_text = shape.chart.chart_title.text_frame.text.strip()
            if title_text:
                segments.append(TextSegment(
                    original_text=title_text,
                    slide_index=slide_idx,
                    shape_index=shape_idx,
                    element_type='chart'
                ))

    return segments


def _extract_from_text_frame(shape: BaseShape, slide_idx: int, shape_idx: int) -> List[TextSegment]:
    """
    テキストフレームからテキストセグメントを抽出する

    Args:
        shape: テキストフレームを持つシェイプ
        slide_idx: スライドインデックス
        shape_idx: シェイプインデックス

    Returns:
        TextSegment のリスト
    """
    segments: List[TextSegment] = []

    # シェイプがプレースホルダーかどうかを判定
    is_title = False
    if hasattr(shape, 'is_placeholder') and shape.is_placeholder:
        try:
            # プレースホルダーの場合のみ placeholder_format にアクセス
            if shape.placeholder_format.type in [1, 3]:  # TITLE or CENTER_TITLE
                is_title = True
        except (AttributeError, ValueError):
            # プレースホルダーでない場合や型が取得できない場合
            pass

    element_type = 'title' if is_title else 'body'

    # パラグラフ単位でテキストを抽出
    for para_idx, paragraph in enumerate(shape.text_frame.paragraphs):
        text = paragraph.text.strip()
        if text:
            segments.append(TextSegment(
                original_text=text,
                slide_index=slide_idx,
                shape_index=shape_idx,
                element_type=element_type,
                paragraph_index=para_idx
            ))

    return segments


def _extract_from_table(shape: BaseShape, slide_idx: int, shape_idx: int) -> List[TextSegment]:
    """
    テーブルからテキストセグメントを抽出する

    Args:
        shape: テーブルシェイプ
        slide_idx: スライドインデックス
        shape_idx: シェイプインデックス

    Returns:
        TextSegment のリスト
    """
    segments: List[TextSegment] = []
    table = shape.table

    for row_idx, row in enumerate(table.rows):
        for col_idx, cell in enumerate(row.cells):
            text = cell.text.strip()
            if text:
                segments.append(TextSegment(
                    original_text=text,
                    slide_index=slide_idx,
                    shape_index=shape_idx,
                    element_type='table',
                    cell_row=row_idx,
                    cell_col=col_idx
                ))

    return segments
