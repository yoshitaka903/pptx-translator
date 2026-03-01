"""
text_extractor.py のテスト
"""

import pytest
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE

from translator.text_extractor import extract_texts, _extract_from_shape, _extract_from_text_frame, _extract_from_table
from translator.models import TextSegment


@pytest.fixture
def empty_presentation():
    """空のプレゼンテーションを作成する"""
    return Presentation()


@pytest.fixture
def simple_presentation():
    """タイトルと本文を持つシンプルなプレゼンテーションを作成する"""
    prs = Presentation()

    # スライド1: タイトルスライド
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])  # Title Slide layout
    slide1.shapes.title.text = "Test Title"
    slide1.placeholders[1].text = "Test Subtitle"

    # スライド2: タイトルと本文
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content layout
    slide2.shapes.title.text = "Slide 2 Title"

    # 本文に箇条書きを追加
    body_shape = slide2.placeholders[1]
    tf = body_shape.text_frame
    tf.text = "First bullet"

    p = tf.add_paragraph()
    p.text = "Second bullet"

    return prs


@pytest.fixture
def presentation_with_table():
    """テーブルを含むプレゼンテーションを作成する"""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout

    # 2x2 のテーブルを追加
    rows, cols = 2, 2
    left = Inches(1.0)
    top = Inches(1.0)
    width = Inches(4.0)
    height = Inches(2.0)

    table = slide.shapes.add_table(rows, cols, left, top, width, height).table

    # セルにテキストを設定
    table.cell(0, 0).text = "Header 1"
    table.cell(0, 1).text = "Header 2"
    table.cell(1, 0).text = "Cell A"
    table.cell(1, 1).text = "Cell B"

    return prs


@pytest.fixture
def presentation_with_notes():
    """スピーカーノートを含むプレゼンテーションを作成する"""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Slide with Notes"

    # ノートを追加
    notes_slide = slide.notes_slide
    text_frame = notes_slide.notes_text_frame
    text_frame.text = "This is a speaker note."

    return prs


class TestExtractTexts:
    """extract_texts 関数のテスト"""

    def test_empty_presentation(self, empty_presentation):
        """空のプレゼンテーションからは何も抽出されない"""
        segments = extract_texts(empty_presentation)
        assert len(segments) == 0

    def test_simple_presentation(self, simple_presentation):
        """シンプルなプレゼンテーションからテキストを抽出できる"""
        segments = extract_texts(simple_presentation, translate_notes=False)

        # スライド1: タイトル + サブタイトル
        # スライド2: タイトル + 箇条書き2つ
        assert len(segments) == 5

        # スライド1のタイトル
        assert segments[0].original_text == "Test Title"
        assert segments[0].slide_index == 0
        assert segments[0].element_type == "title"

        # スライド1のサブタイトル
        assert segments[1].original_text == "Test Subtitle"
        assert segments[1].slide_index == 0
        assert segments[1].element_type == "body"

        # スライド2のタイトル
        assert segments[2].original_text == "Slide 2 Title"
        assert segments[2].slide_index == 1
        assert segments[2].element_type == "title"

    def test_table_extraction(self, presentation_with_table):
        """テーブルからテキストを抽出できる"""
        segments = extract_texts(presentation_with_table)

        assert len(segments) == 4

        # すべてテーブル要素
        for seg in segments:
            assert seg.element_type == "table"
            assert seg.cell_row is not None
            assert seg.cell_col is not None

        # Header 1 をチェック
        header1 = [s for s in segments if s.original_text == "Header 1"][0]
        assert header1.cell_row == 0
        assert header1.cell_col == 0

    def test_notes_extraction(self, presentation_with_notes):
        """スピーカーノートを抽出できる"""
        segments = extract_texts(presentation_with_notes, translate_notes=True)

        # タイトル + ノート
        assert len(segments) == 2

        note_segment = segments[1]
        assert note_segment.original_text == "This is a speaker note."
        assert note_segment.element_type == "note"
        assert note_segment.shape_index == -1

    def test_notes_skip(self, presentation_with_notes):
        """translate_notes=False でノートをスキップできる"""
        segments = extract_texts(presentation_with_notes, translate_notes=False)

        # タイトルのみ
        assert len(segments) == 1
        assert segments[0].element_type == "title"


class TestExtractFromShape:
    """_extract_from_shape 関数のテスト"""

    def test_text_frame_shape(self, simple_presentation):
        """テキストフレームを持つシェイプから抽出できる"""
        slide = simple_presentation.slides[0]
        title_shape = slide.shapes.title

        segments = _extract_from_shape(title_shape, 0, 0)

        assert len(segments) == 1
        assert segments[0].original_text == "Test Title"
        assert segments[0].element_type == "title"

    def test_empty_text_frame(self):
        """空のテキストフレームは無視される"""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        # タイトルは空のまま

        segments = _extract_from_shape(slide.shapes.title, 0, 0)
        assert len(segments) == 0

    def test_table_shape(self, presentation_with_table):
        """テーブルシェイプから抽出できる"""
        slide = presentation_with_table.slides[0]
        # テーブルは GraphicFrame オブジェクトなので、has_table でフィルタ
        table_shape = [s for s in slide.shapes if hasattr(s, 'has_table') and s.has_table][0]

        segments = _extract_from_shape(table_shape, 0, 0)

        assert len(segments) == 4
        assert all(s.element_type == "table" for s in segments)


class TestExtractFromTextFrame:
    """_extract_from_text_frame 関数のテスト"""

    def test_title_detection(self, simple_presentation):
        """タイトルプレースホルダーを正しく検出できる"""
        slide = simple_presentation.slides[0]
        title_shape = slide.shapes.title

        segments = _extract_from_text_frame(title_shape, 0, 0)

        assert len(segments) == 1
        assert segments[0].element_type == "title"

    def test_body_detection(self, simple_presentation):
        """本文プレースホルダーを正しく検出できる"""
        slide = simple_presentation.slides[0]
        subtitle_shape = slide.placeholders[1]

        segments = _extract_from_text_frame(subtitle_shape, 0, 1)

        assert len(segments) == 1
        assert segments[0].element_type == "body"

    def test_multiple_paragraphs(self, simple_presentation):
        """複数のパラグラフを個別に抽出できる"""
        slide = simple_presentation.slides[1]
        body_shape = slide.placeholders[1]

        segments = _extract_from_text_frame(body_shape, 1, 1)

        assert len(segments) == 2
        assert segments[0].original_text == "First bullet"
        assert segments[0].paragraph_index == 0
        assert segments[1].original_text == "Second bullet"
        assert segments[1].paragraph_index == 1


class TestExtractFromTable:
    """_extract_from_table 関数のテスト"""

    def test_table_cells(self, presentation_with_table):
        """テーブルのすべてのセルから抽出できる"""
        slide = presentation_with_table.slides[0]
        # テーブルは GraphicFrame オブジェクトなので、has_table でフィルタ
        table_shape = [s for s in slide.shapes if hasattr(s, 'has_table') and s.has_table][0]

        segments = _extract_from_table(table_shape, 0, 0)

        assert len(segments) == 4

        # 座標のチェック
        cell_coords = {(s.cell_row, s.cell_col) for s in segments}
        expected_coords = {(0, 0), (0, 1), (1, 0), (1, 1)}
        assert cell_coords == expected_coords

    def test_empty_cells_skipped(self):
        """空のセルはスキップされる"""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])

        table = slide.shapes.add_table(2, 2, Inches(1), Inches(1), Inches(4), Inches(2)).table
        table.cell(0, 0).text = "Only this cell has text"
        # 他のセルは空のまま

        # テーブルは GraphicFrame オブジェクトなので、has_table でフィルタ
        table_shape = [s for s in slide.shapes if hasattr(s, 'has_table') and s.has_table][0]
        segments = _extract_from_table(table_shape, 0, 0)

        assert len(segments) == 1
        assert segments[0].original_text == "Only this cell has text"
        assert segments[0].cell_row == 0
        assert segments[0].cell_col == 0
