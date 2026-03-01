"""
pptx_handler.py のテスト
"""

import pytest
from pptx import Presentation
from pptx.util import Inches

from translator.pptx_handler import apply_translations, _update_paragraph_text
from translator.models import TextSegment


@pytest.fixture
def simple_presentation():
    """タイトルと本文を持つシンプルなプレゼンテーションを作成する"""
    prs = Presentation()

    # スライド1: タイトルスライド
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])  # Title Slide layout
    slide1.shapes.title.text = "Test Title"
    slide1.placeholders[1].text = "Test Subtitle"

    # スライド2: タイトルと本文（複数パラグラフ）
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content layout
    slide2.shapes.title.text = "Slide 2 Title"

    body_shape = slide2.placeholders[1]
    tf = body_shape.text_frame
    tf.text = "First bullet"

    p = tf.add_paragraph()
    p.text = "Second bullet"

    return prs


@pytest.fixture
def presentation_with_table():
    """テーブルを含むプレゼンテーションを作成する"""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout

    # 2x2 のテーブルを追加
    rows, cols = 2, 2
    left = Inches(1.0)
    top = Inches(1.0)
    width = Inches(4.0)
    height = Inches(2.0)

    table = slide.shapes.add_table(rows, cols, left, top, width, height).table

    # セルにテキストを設定
    table.cell(0, 0).text = "Header 1"
    table.cell(0, 1).text = "Header 2"
    table.cell(1, 0).text = "Cell A"
    table.cell(1, 1).text = "Cell B"

    return prs


@pytest.fixture
def presentation_with_notes():
    """スピーカーノートを含むプレゼンテーションを作成する"""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Slide with Notes"

    # ノートを追加
    notes_slide = slide.notes_slide
    text_frame = notes_slide.notes_text_frame
    text_frame.text = "This is a speaker note."

    return prs


class TestApplyTranslations:
    """apply_translations 関数のテスト"""

    def test_empty_segments(self, simple_presentation):
        """空のセグメントリストの場合は何も変更しない"""
        original_title = simple_presentation.slides[0].shapes.title.text

        result = apply_translations(simple_presentation, [])

        # 元のテキストが変更されていないことを確認
        assert result.slides[0].shapes.title.text == original_title

    def test_apply_to_title(self, simple_presentation):
        """タイトルに翻訳を適用できる"""
        segments = [
            TextSegment(
                original_text="Test Title",
                slide_index=0,
                shape_index=0,
                element_type="title",
                translated_text="テストタイトル",
            )
        ]

        result = apply_translations(simple_presentation, segments)

        # 翻訳が適用されていることを確認
        assert result.slides[0].shapes.title.text == "テストタイトル"

    def test_apply_to_multiple_paragraphs(self, simple_presentation):
        """複数のパラグラフに個別に翻訳を適用できる"""
        segments = [
            TextSegment(
                original_text="First bullet",
                slide_index=1,
                shape_index=1,
                element_type="body",
                translated_text="最初の箇条書き",
                paragraph_index=0,
            ),
            TextSegment(
                original_text="Second bullet",
                slide_index=1,
                shape_index=1,
                element_type="body",
                translated_text="2番目の箇条書き",
                paragraph_index=1,
            ),
        ]

        result = apply_translations(simple_presentation, segments)

        # 各パラグラフの翻訳を確認
        body_shape = result.slides[1].placeholders[1]
        paragraphs = body_shape.text_frame.paragraphs

        assert paragraphs[0].text == "最初の箇条書き"
        assert paragraphs[1].text == "2番目の箇条書き"

    def test_skip_segments_without_translation(self, simple_presentation):
        """translated_text が None のセグメントはスキップする"""
        original_title = simple_presentation.slides[0].shapes.title.text

        segments = [
            TextSegment(
                original_text="Test Title",
                slide_index=0,
                shape_index=0,
                element_type="title",
                translated_text=None,  # 翻訳なし
            )
        ]

        result = apply_translations(simple_presentation, segments)

        # 元のテキストが変更されていないことを確認
        assert result.slides[0].shapes.title.text == original_title

    def test_apply_to_table(self, presentation_with_table):
        """テーブルセルに翻訳を適用できる"""
        # テーブルは GraphicFrame オブジェクトなので、has_table でフィルタ
        slide = presentation_with_table.slides[0]
        table_shape_index = next(
            i for i, s in enumerate(slide.shapes)
            if hasattr(s, 'has_table') and s.has_table
        )

        segments = [
            TextSegment(
                original_text="Header 1",
                slide_index=0,
                shape_index=table_shape_index,
                element_type="table",
                translated_text="ヘッダー1",
                cell_row=0,
                cell_col=0,
            ),
            TextSegment(
                original_text="Cell A",
                slide_index=0,
                shape_index=table_shape_index,
                element_type="table",
                translated_text="セルA",
                cell_row=1,
                cell_col=0,
            ),
        ]

        result = apply_translations(presentation_with_table, segments)

        # テーブルの翻訳を確認
        table_shape = [s for s in result.slides[0].shapes if hasattr(s, 'has_table') and s.has_table][0]
        table = table_shape.table

        assert table.cell(0, 0).text == "ヘッダー1"
        assert table.cell(1, 0).text == "セルA"

    def test_apply_to_note(self, presentation_with_notes):
        """スピーカーノートに翻訳を適用できる"""
        segments = [
            TextSegment(
                original_text="This is a speaker note.",
                slide_index=0,
                shape_index=-1,
                element_type="note",
                translated_text="これはスピーカーノートです。",
            )
        ]

        result = apply_translations(presentation_with_notes, segments)

        # ノートの翻訳を確認
        notes_text = result.slides[0].notes_slide.notes_text_frame.text
        assert notes_text == "これはスピーカーノートです。"


class TestUpdateParagraphText:
    """_update_paragraph_text 関数のテスト"""

    def test_update_single_run(self):
        """単一 Run のパラグラフを更新できる"""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = "Original Text"

        paragraph = slide.shapes.title.text_frame.paragraphs[0]
        _update_paragraph_text(paragraph, "Updated Text")

        assert paragraph.text == "Updated Text"

    def test_update_multiple_runs_preserves_first_style(self):
        """複数 Run の場合、最初の Run にテキストを設定し、残りは空文字にする"""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])

        # 複数の Run を作成
        title = slide.shapes.title
        tf = title.text_frame
        tf.clear()

        p = tf.paragraphs[0]
        run1 = p.add_run()
        run1.text = "First"
        run2 = p.add_run()
        run2.text = " Second"
        run3 = p.add_run()
        run3.text = " Third"

        # 更新前の Run 数を確認
        assert len(p.runs) == 3

        # 更新
        _update_paragraph_text(p, "New Text")

        # 最初の Run にのみテキストが設定され、残りは空文字
        assert p.runs[0].text == "New Text"
        assert p.runs[1].text == ""
        assert p.runs[2].text == ""

        # Run の数は変わらない
        assert len(p.runs) == 3

    def test_update_empty_paragraph(self):
        """Run がないパラグラフの場合は新規作成される"""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank

        # テキストボックスを追加
        left = Inches(1)
        top = Inches(1)
        width = Inches(3)
        height = Inches(1)
        textbox = slide.shapes.add_textbox(left, top, width, height)

        paragraph = textbox.text_frame.paragraphs[0]

        # 更新前は Run がない
        assert len(paragraph.runs) == 0

        # 更新
        _update_paragraph_text(paragraph, "New Text")

        # テキストが設定される
        assert paragraph.text == "New Text"


class TestStylePreservation:
    """スタイル維持のテスト"""

    def test_font_style_preserved(self):
        """フォントスタイルが維持されることを確認"""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])

        title = slide.shapes.title
        tf = title.text_frame
        tf.text = "Original"

        # フォントスタイルを設定
        run = tf.paragraphs[0].runs[0]
        run.font.bold = True
        run.font.italic = True
        run.font.size = Inches(0.5)

        # 元のスタイルを保存
        original_bold = run.font.bold
        original_italic = run.font.italic
        original_size = run.font.size

        # 翻訳を適用
        segments = [
            TextSegment(
                original_text="Original",
                slide_index=0,
                shape_index=0,
                element_type="title",
                translated_text="翻訳後",
            )
        ]

        apply_translations(prs, segments)

        # スタイルが維持されていることを確認
        updated_run = tf.paragraphs[0].runs[0]
        assert updated_run.text == "翻訳後"
        assert updated_run.font.bold == original_bold
        assert updated_run.font.italic == original_italic
        assert updated_run.font.size == original_size
